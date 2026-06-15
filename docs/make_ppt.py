from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

BRAND_BLUE = RGBColor(0x1d, 0x4e, 0xd8)
LIGHT_BLUE = RGBColor(0xdb, 0xe9, 0xfe)
DARK = RGBColor(0x1e, 0x29, 0x3b)
WHITE = RGBColor(0xff, 0xff, 0xff)
GRAY = RGBColor(0x64, 0x74, 0x8b)

blank = prs.slide_layouts[6]  # blank

def add_rect(slide, l, t, w, h, color):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h, size=18, bold=False, color=RGBColor(0x1e,0x29,0x3b), align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.size = Pt(size); run.font.bold = bold
    run.font.color.rgb = color; run.font.italic = italic
    return tb

# ── 슬라이드 1: 표지 ──────────────────────────────────────────
s1 = prs.slides.add_slide(blank)
add_rect(s1, 0, 0, 13.33, 7.5, BRAND_BLUE)
add_rect(s1, 0, 4.8, 13.33, 2.7, RGBColor(0x17,0x3a,0xa8))
add_text(s1, '더원컴퍼니', 1, 1.2, 11, 1, size=24, bold=True, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
add_text(s1, '업무관리시스템', 1, 2.0, 11, 1.2, size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s1, '사용자 매뉴얼', 1, 3.1, 11, 0.8, size=22, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
add_text(s1, '2026-03-23  |  http://43.203.193.217:10000', 1, 5.2, 11, 0.6, size=13, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s1, '대전광역시 둔산동 1229 6층  |  02-1660-4130', 1, 5.8, 11, 0.6, size=12, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

# ── 슬라이드 2: 목차 ──────────────────────────────────────────
s2 = prs.slides.add_slide(blank)
add_rect(s2, 0, 0, 13.33, 1.2, BRAND_BLUE)
add_text(s2, '목  차', 0.5, 0.2, 12, 0.8, size=26, bold=True, color=WHITE)
items = ['01  시스템 개요','02  권한 체계','03  로그인 및 로그아웃',
         '04  비밀번호 관리','05  공지사항','06  고객관리',
         '07  인사관리 (관리자 전용)','08  공통코드관리','09  세션 관리 / FAQ']
cols = [items[:5], items[5:]]
for ci, col in enumerate(cols):
    for ri, item in enumerate(col):
        add_rect(s2, 0.5+ci*6.4, 1.5+ri*1.0, 5.8, 0.8, LIGHT_BLUE)
        add_text(s2, item, 0.7+ci*6.4, 1.55+ri*1.0, 5.4, 0.7, size=15, bold=True, color=BRAND_BLUE)

# ── 헬퍼: 섹션 슬라이드 ──────────────────────────────────────
def section_slide(title, items_list, note=''):
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, 13.33, 1.2, BRAND_BLUE)
    add_text(s, title, 0.5, 0.2, 12, 0.8, size=24, bold=True, color=WHITE)
    y = 1.4
    for item in items_list:
        if item.startswith('##'):
            add_text(s, item[2:].strip(), 0.5, y, 12, 0.5, size=14, bold=True, color=BRAND_BLUE)
            y += 0.5
        elif item.startswith('•'):
            add_rect(s, 0.5, y+0.12, 0.08, 0.08, BRAND_BLUE)
            add_text(s, item[1:].strip(), 0.7, y, 11.8, 0.45, size=13, color=DARK)
            y += 0.48
        else:
            add_text(s, item, 0.5, y, 12, 0.45, size=13, color=DARK)
            y += 0.48
    if note:
        add_rect(s, 0.5, 6.8, 12.3, 0.5, LIGHT_BLUE)
        add_text(s, '※ ' + note, 0.6, 6.82, 12, 0.4, size=11, color=BRAND_BLUE, italic=True)
    return s

# ── 슬라이드 3: 시스템 개요 ───────────────────────────────────
section_slide('01  시스템 개요', [
    '더원컴퍼니 내부 업무 관리 시스템입니다.',
    '',
    '## 주요 기능',
    '• 공지사항 조회 및 작성 (권한에 따라)',
    '• 개인고객 및 협력점 고객 정보 관리',
    '• 직원 계정 관리 및 권한 설정 (관리자)',
    '• 공통코드 관리 (드롭다운 항목 관리)',
    '',
    '## 접속 정보',
    '• 시스템 주소: http://43.203.193.217:10000',
    '• 권장 브라우저: Chrome, Edge (최신 버전)',
])

# ── 슬라이드 4: 권한 체계 ─────────────────────────────────────
s4 = prs.slides.add_slide(blank)
add_rect(s4, 0, 0, 13.33, 1.2, BRAND_BLUE)
add_text(s4, '02  권한 체계', 0.5, 0.2, 12, 0.8, size=24, bold=True, color=WHITE)

roles_data = [
    ('ADMIN', '최고관리자', ['공지사항 작성/수정/삭제', '고객관리', '인사관리', '공통코드관리'], BRAND_BLUE),
    ('MANAGER', '매니저', ['공지사항 작성/수정/삭제', '고객관리', '공통코드관리'], RGBColor(0x06,0x95,0x6b)),
    ('MEMBER', '일반사용자', ['공지사항 읽기만 가능', '고객관리'], RGBColor(0x97,0x16,0x16)),
]
for ci, (role, name, perms, color) in enumerate(roles_data):
    x = 0.4 + ci * 4.3
    add_rect(s4, x, 1.3, 4.0, 0.7, color)
    add_text(s4, role, x, 1.35, 4.0, 0.6, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s4, name, x, 2.0, 4.0, 0.5, size=13, color=color, align=PP_ALIGN.CENTER, bold=True)
    for pi, perm in enumerate(perms):
        add_rect(s4, x+0.1, 2.7+pi*0.65, 0.08, 0.08, color)
        add_text(s4, perm, x+0.3, 2.62+pi*0.65, 3.6, 0.55, size=12, color=DARK)

add_rect(s4, 0.4, 6.6, 12.5, 0.5, LIGHT_BLUE)
add_text(s4, '※ 권한이 없는 메뉴 접근 시 공지사항 화면으로 자동 이동됩니다.', 0.6, 6.62, 12, 0.4, size=11, color=BRAND_BLUE, italic=True)

# ── 슬라이드 5: 로그인 ────────────────────────────────────────
section_slide('03  로그인 및 로그아웃', [
    '## 로그인 방법',
    '• 브라우저에서 시스템 주소 접속 → 자동으로 로그인 화면 표시',
    '• 아이디와 비밀번호 입력 후 로그인 버튼 클릭',
    '• 로그인 성공 시 공지사항 화면으로 이동',
    '',
    '## 오류 안내',
    '• 아이디/비밀번호 오류 → "아이디 또는 비밀번호가 올바르지 않습니다" (401)',
    '• 비활성 계정 → "비활성화된 계정입니다" → 관리자에게 문의',
    '• 비밀번호 초기화된 계정 → 로그인 즉시 재설정 화면으로 이동',
    '',
    '## 로그아웃',
    '• 우측 상단 메뉴에서 로그아웃 클릭 → 로그인 이력 자동 저장',
])

# ── 슬라이드 6: 비밀번호 관리 ────────────────────────────────
section_slide('04  비밀번호 관리', [
    '## 비밀번호 초기화 후 재설정 (관리자가 초기화한 경우)',
    '• 로그인 시 자동으로 재설정 화면 표시',
    '• 새 비밀번호 입력 → 확인 입력 → 변경 버튼 클릭',
    '• 초기 비밀번호: 아이디 + 1234!  (예: hong → hong1234!)',
    '',
    '## 내 비밀번호 변경 (자발적 변경)',
    '• 우측 상단 메뉴에서 비밀번호 변경 선택',
    '• 현재 비밀번호 입력 → 새 비밀번호 입력 → 저장',
    '',
    '## 비밀번호 규칙',
    '• 8자 이상',
    '• 영문 + 숫자 + 특수문자를 모두 포함',
    '• 예시: MyPass1!  /  Work2026@',
])

# ── 슬라이드 7: 공지사항 ──────────────────────────────────────
section_slide('05  공지사항', [
    '로그인 후 기본으로 표시되는 화면입니다.',
    '',
    '## 권한별 기능',
    '• ADMIN / MANAGER: 공지사항 작성, 수정(본인 글), 삭제(본인 글)',
    '• ADMIN: 모든 공지사항 삭제 가능',
    '• MEMBER: 공지사항 읽기만 가능 (작성 버튼 미표시)',
    '',
    '## 사용 방법',
    '• 목록에서 제목 클릭 → 상세 내용 확인',
    '• 작성 버튼(+) 클릭 → 제목/내용 입력 → 저장',
    '• 수정: 목록의 수정 아이콘 클릭 → 내용 수정 → 저장',
], note='MEMBER 권한은 공지사항 읽기 전용입니다.')

# ── 슬라이드 8: 고객관리 ──────────────────────────────────────
section_slide('06  고객관리', [
    '## 개인고객관리',
    '• 조회: 고객명, 연락처, 상태 등 조건 입력 후 조회 버튼 클릭',
    '• 등록: 등록 버튼 클릭 → 고객 정보 입력 → 저장',
    '• 수정: 목록에서 고객 선택 → 정보 수정 → 저장',
    '• 삭제: 목록에서 고객 선택 → 삭제 버튼 클릭',
    '',
    '## 협력점 고객관리',
    '• 협력점(파트너사) 고객 정보 관리',
    '• 조작 방법은 개인고객관리와 동일',
])

# ── 슬라이드 9: 인사관리 ──────────────────────────────────────
section_slide('07  인사관리 (ADMIN 전용)', [
    '## 직원 등록',
    '• 등록 버튼 클릭 → 아이디, 비밀번호, 이름, 부서, 권한, 사용여부 입력 → 저장',
    '• 아이디 중복 확인 필수',
    '',
    '## 직원 수정',
    '• 목록 선택 → 정보 수정 → 저장',
    '• 비밀번호는 새로 입력할 때만 변경됨 (공란이면 기존 유지)',
    '',
    '## 비밀번호 초기화',
    '• 초기화 버튼 클릭 → 아이디+1234! 로 설정',
    '• 해당 직원은 다음 로그인 시 반드시 비밀번호 변경 필요',
    '',
    '## 직원 삭제',
    '• 목록 선택 → 삭제 (로그인 이력 포함 삭제)',
], note='인사관리는 ADMIN 권한만 접근 가능합니다.')

# ── 슬라이드 10: 공통코드 + 세션 ─────────────────────────────
section_slide('08  공통코드관리 / 09  세션 관리', [
    '## 공통코드관리 (ADMIN, MANAGER)',
    '• 시스템 드롭다운 선택 항목(코드 값)을 관리',
    '• 그룹코드 관리: 코드 그룹 단위 조회/등록/수정/삭제',
    '• 상세코드 관리: 그룹 내 개별 코드 항목 관리',
    '• 사용여부 N 설정 시 드롭다운 목록에서 숨겨짐',
    '',
    '## 세션 관리',
    '• 로그인 후 1시간 경과 시 자동 로그아웃',
    '• 만료 시 "세션이 만료되었습니다. 다시 로그인해주세요." 메시지 표시',
    '• 여러 탭 사용 시 한 탭 로그아웃 → 모든 탭 동시 로그아웃',
])

# ── 슬라이드 11: FAQ ──────────────────────────────────────────
s11 = prs.slides.add_slide(blank)
add_rect(s11, 0, 0, 13.33, 1.2, BRAND_BLUE)
add_text(s11, '자주 묻는 질문 (FAQ)', 0.5, 0.2, 12, 0.8, size=24, bold=True, color=WHITE)
faqs = [
    ('비밀번호를 잊어버렸습니다.', '관리자에게 초기화 요청 → 아이디+1234! 로 로그인 → 비밀번호 변경'),
    ('메뉴가 보이지 않습니다.', '권한 부족입니다. 관리자에게 권한 변경 요청하세요.'),
    ('로그인이 안 됩니다.', '계정이 비활성 상태일 수 있습니다. 관리자에게 활성화 요청하세요.'),
    ('갑자기 로그인 화면으로 이동합니다.', '세션(1시간)이 만료되었습니다. 다시 로그인하세요.'),
    ('공지사항 작성 버튼이 없습니다.', 'MEMBER 권한은 읽기 전용입니다. 작성은 ADMIN/MANAGER만 가능합니다.'),
]
for i, (q, a) in enumerate(faqs):
    y = 1.4 + i * 1.0
    add_rect(s11, 0.4, y, 12.5, 0.85, LIGHT_BLUE)
    add_text(s11, 'Q.  ' + q, 0.6, y+0.05, 12, 0.38, size=13, bold=True, color=BRAND_BLUE)
    add_text(s11, 'A.  ' + a, 0.6, y+0.43, 12, 0.38, size=12, color=DARK)

# ── 슬라이드 12: 마무리 ───────────────────────────────────────
s12 = prs.slides.add_slide(blank)
add_rect(s12, 0, 0, 13.33, 7.5, BRAND_BLUE)
add_rect(s12, 0, 3.0, 13.33, 1.5, RGBColor(0x17,0x3a,0xa8))
add_text(s12, '감사합니다', 1, 1.5, 11, 1.5, size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s12, '문의사항이 있으시면 관리자에게 연락해주세요.', 1, 3.1, 11, 0.7, size=16, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
add_text(s12, '더원컴퍼니  |  대전광역시 둔산동 1229 6층  |  02-1660-4130', 1, 5.5, 11, 0.6, size=14, color=WHITE, align=PP_ALIGN.CENTER)

out = 'C:/Users/ocean-nopc/Desktop/project/managesys/docs/더원컴퍼니_업무관리시스템_사용자매뉴얼.pptx'
prs.save(out)
print('PPT 저장 완료:', out)
